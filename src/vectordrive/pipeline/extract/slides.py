"""PPTX extraction via python-pptx: one slide becomes one ExtractedPage
(page_number = slide's 1-indexed position in the deck), same reasoning as
spreadsheet.py's one-sheet-per-page choice — reusing the existing
page/citation model means chunking, embedding, search, and MCP citations
need zero changes to support a new document type.

A slide's text is every shape's text frame (title, body, any text box), in
shape order, followed by its speaker notes if present — notes often carry
the actual explanatory content behind a terse slide, and are otherwise
invisible to any search over the deck.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from vectordrive.pipeline.extract.base import ExtractedDocument, ExtractedPage


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            parts.append(shape.text_frame.text)
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            parts.append(f"Notes: {notes}")
    return "\n".join(parts)


class SlideExtractor:
    NAME = "python-pptx"
    VERSION = "1"

    def extract(self, path: Path) -> ExtractedDocument:
        prs = Presentation(path)
        pages = [
            ExtractedPage(
                page_number=index, text=_slide_text(slide), text_source="native", needs_ocr=False
            )
            for index, slide in enumerate(prs.slides, start=1)
        ]
        return ExtractedDocument(
            extractor_name=self.NAME,
            extractor_version=self.VERSION,
            page_count=len(pages),
            pages=pages,
        )
